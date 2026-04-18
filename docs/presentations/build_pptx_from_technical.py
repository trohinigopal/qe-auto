"""
Generate qe-auto-technical.pptx from the same narrative as docs/TECHNICAL.md.

Run from repo root or framework (paths resolve relative to this file):
  pip install python-pptx
  python docs/presentations/build_pptx_from_technical.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[2]
OUT = Path(__file__).resolve().parent / "qe-auto-technical.pptx"

ACCENT = RGBColor(34, 211, 238)
MUTED = RGBColor(100, 120, 140)


def _add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    sub = slide.placeholders[1]
    sub.text = subtitle


def _add_bullets(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(18)
        if p.font.color:
            p.font.color.rgb = RGBColor(30, 41, 59)


def _add_table_slide(prs: Presentation, title: str, headers: list[str], rows: list[list[str]]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_box.text_frame.text = title
    title_box.text_frame.paragraphs[0].font.size = Pt(28)
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].font.color.rgb = ACCENT

    rows_n = len(rows) + 1
    cols_n = len(headers)
    table = slide.shapes.add_table(
        Inches(0.5), Inches(1.3), Inches(9), Inches(4.5), rows_n, cols_n
    ).table
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            table.cell(r, c).text = val
            table.cell(r, c).text_frame.paragraphs[0].font.size = Pt(12)


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _add_title_slide(
        prs,
        "qe-auto — Technical overview",
        "Multi-channel BDD automation · pytest · Selenium/Playwright · Appium · API · Allure\n"
        + f"Source: {ROOT / 'docs' / 'TECHNICAL.md'}",
    )

    _add_bullets(
        prs,
        "Purpose & channels",
        [
            "Netbanking — Selenium/Playwright (web UI)",
            "MobileBanking — Appium",
            "API — requests via ApiBaseClient",
            "CrossChannel — web + API in one journey",
            "Gherkin (.feature) + pytest-bdd; fixtures: ui_driver, mobile_driver, api_client",
        ],
    )

    _add_bullets(
        prs,
        "Repository layout",
        [
            "framework/ — runner.py, conftest.py, pytest.ini, requirements.txt",
            "Per channel: Features/nm/, pageobject/nm/, stepdefination/nm/",
            "common/: base/, reporting/, utils/, steps/",
            "Same <stem> for feature, page object, and test_*_steps.py",
        ],
    )

    _add_bullets(
        prs,
        "Prerequisites",
        [
            "Python 3.10+ · pip · Chrome (+ driver) · Git",
            "Optional: Allure CLI (npm), Playwright browsers, Appium + SDKs, Firefox",
        ],
    )

    _add_table_slide(
        prs,
        "CLI options (conftest)",
        ["Option", "Default", "Role"],
        [
            ["--platform", "markers/web", "web | mobile | api | cross"],
            ["--browser", "chrome", "Selenium browser"],
            ["--web-engine", "selenium", "selenium | playwright"],
            ["--base-url", "example-bank.test", "Web SUT"],
            ["--api-base-url", "example-bank-api.test", "API base"],
        ],
    )

    _add_bullets(
        prs,
        "pytest_plugins & step_defs",
        [
            "Step code in *_step_defs.py must be loaded via pytest_plugins",
            "Netbanking/stepdefination/conftest.py → NB step_defs",
            "API/stepdefination/conftest.py → API step_defs",
            "CrossChannel/conftest.py → both for cross scenarios",
            "Otherwise: StepDefinitionNotFoundError",
        ],
    )

    _add_bullets(
        prs,
        "Generators & materialize",
        [
            "user_story_to_manual — UserStories → ManualTestCases",
            "manual_tc_to_bdd — Manual TC → .feature + page + steps",
            "bdd_dom_materialize — DOM scan → *_dom_materialized.json + patch steps",
            "api_bdd_materialize — *_api_materialized.json + API step runner",
        ],
    )

    _add_bullets(
        prs,
        "Allure reporting",
        [
            "Session folder: common/reports/sessions/<id>_<timestamp>/",
            "export_allure_artifacts → HTML (needs CLI), xlsx, pdf, report_hub.html",
            "Screenshots on fail; Playwright video if ALLURE_RECORD_VIDEO=1",
        ],
    )

    _add_bullets(
        prs,
        "Runner commands",
        [
            "python runner.py --platform cross|web|api|mobile",
            "python runner.py --tags \"@smoke\" --tests-path <file>",
            "python runner.py --build-allure [--open-allure]",
            "python runner.py --quality-check  ·  --build-rtm",
        ],
    )

    _add_bullets(
        prs,
        "Troubleshooting",
        [
            "Step not found → check stepdefination/conftest.py pytest_plugins",
            "No Allure HTML → install allure CLI",
            "API issues → --platform api or cross; real URLs",
        ],
    )

    _add_bullets(
        prs,
        "Documentation",
        [
            "docs/TECHNICAL.md — full reference",
            "framework/README.md — quick commands",
            "docs/presentations/qe-auto-technical-deck.html — HTML deck",
        ],
    )

    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    p = build()
    print(f"Wrote: {p}")
